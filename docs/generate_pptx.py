"""
Script de generation du PowerPoint - BookingSysteme
Presentation de 10-15 minutes couvrant tous les aspects du projet.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Couleurs du theme ---
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Fond sombre
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xC7)    # Bleu accent
ACCENT_ORANGE = RGBColor(0xF7, 0x7F, 0x00)  # Orange accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x88)
SECTION_BG = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_dark_bg(slide, color=DARK_BG):
    """Ajoute un fond sombre a un slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Ajoute un rectangle colore comme fond partiel."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Ajoute un textbox avec du texte formate."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT_BLUE):
    """Ajoute une liste a puces."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Segoe UI"

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Segoe UI"
    return txBox


def add_card(slide, left, top, width, height, title, items, accent_color=ACCENT_BLUE):
    """Ajoute une carte avec titre et contenu."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=15, color=accent_color, bold=True)

    # Items
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4),
                    height - Inches(0.7), items, font_size=12, color=LIGHT_GRAY, bullet_color=accent_color)


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Ajoute une ligne decorative."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_section_divider(title, subtitle=""):
    """Cree un slide de separation de section."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_bg(slide, SECTION_BG)

    # Grande ligne decorative
    add_shape_bg(slide, Inches(0), Inches(3.2), W, Inches(0.06), ACCENT_BLUE)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
                     subtitle, font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


# =====================================================================
# SLIDE 1 - PAGE DE TITRE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Bande decorative en haut
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(0.08), W, Inches(0.04), ACCENT_ORANGE)

# Titre principal
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0),
             "BookingSysteme", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Sous-titre
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.2),
             "Plateforme de reservation multi-services\nRestauration | Hebergement | Transport | Prestations",
             font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(4.0), Inches(4.3), ACCENT_ORANGE)

add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.8),
             "Allocation multi-ressources, optimisation d'occupation\net disponibilite en temps reel",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Contexte
add_text_box(slide, Inches(3), Inches(5.8), Inches(7.3), Inches(0.5),
             "Adapte au contexte Africain | Equipe PING 8",
             font_size=16, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(6.5), Inches(7.3), Inches(0.5),
             "Fevrier 2026",
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2 - SOMMAIRE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Sommaire", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

sommaire_items = [
    ("01", "Contexte & Problematique"),
    ("02", "Objectifs du projet"),
    ("03", "Architecture technique"),
    ("04", "Les microservices"),
    ("05", "Specificites contexte africain"),
    ("06", "Stack technique & outils"),
    ("07", "Maquettes & UX"),
    ("08", "Organisation & methodologie"),
    ("09", "Etat d'avancement"),
    ("10", "Demo & perspectives"),
]

for i, (num, label) in enumerate(sommaire_items):
    row = i // 2
    col = i % 2
    x = Inches(1.0) + col * Inches(6)
    y = Inches(1.6) + row * Inches(1.0)

    add_text_box(slide, x, y, Inches(0.6), Inches(0.5),
                 num, font_size=22, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.03), Inches(5), Inches(0.5),
                 label, font_size=18, color=WHITE)


# =====================================================================
# SLIDE 3 - CONTEXTE & PROBLEMATIQUE
# =====================================================================
add_section_divider("01 - Contexte & Problematique",
                    "Pourquoi ce projet ?")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Contexte & Problematique", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3.5))

# Colonne gauche - Constat
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3),
         "CONSTAT",
         [
             "Gestion manuelle des reservations (papier, WhatsApp, telephone)",
             "Conflits frequents : doubles reservations, oublis",
             "Aucune visibilite en temps reel sur les disponibilites",
             "Perte de revenus : ressources sous-utilisees",
             "Pas d'outil adapte au marche africain",
             "90% des paiements se font par Mobile Money",
         ], accent_color=RED)

# Colonne droite - Solution
add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3),
         "NOTRE SOLUTION",
         [
             "Plateforme unifiee de reservation multi-services",
             "Disponibilite en temps reel (WebSocket)",
             "Allocation intelligente des ressources",
             "Paiement Mobile Money integre (Orange, MTN, Wave...)",
             "Paiement en especes sur place (contexte local)",
             "Notifications automatiques (confirmation, rappels)",
         ], accent_color=GREEN)


# =====================================================================
# SLIDE 4 - OBJECTIFS
# =====================================================================
add_section_divider("02 - Objectifs du projet",
                    "Ce que nous voulons atteindre")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Objectifs du projet", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3))

# 4 cartes d'objectifs
objectives = [
    ("Pour les Commercants", ACCENT_BLUE, [
        "Gerer tables, chambres, bus, evenements",
        "Eviter les conflits de reservation",
        "Maximiser le taux d'occupation",
        "Dashboard de suivi des revenus",
    ]),
    ("Pour les Clients", GREEN, [
        "Reserver en ligne facilement",
        "Voir les disponibilites en temps reel",
        "Suggestions de creneaux alternatifs",
        "Payer via Mobile Money ou carte",
    ]),
    ("Pour la Plateforme", ACCENT_ORANGE, [
        "Modele economique par commission",
        "Multi-services : restaurant, hotel, transport...",
        "Scalable et containerise",
        "Documentation API complete",
    ]),
    ("Contexte Africain", RGBColor(0xAF, 0x7A, 0xC5), [
        "Mobile Money comme moyen principal",
        "Support paiement en especes sur place",
        "Devises FCFA (XOF, XAF), USD, EUR",
        "Evenements familiaux (mariage, bapteme)",
    ]),
]

for i, (title, color, items) in enumerate(objectives):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.3)
    y = Inches(1.5) + row * Inches(2.85)
    add_card(slide, x, y, Inches(5.9), Inches(2.6), title, items, accent_color=color)


# =====================================================================
# SLIDE 5 - ARCHITECTURE TECHNIQUE (VUE D'ENSEMBLE)
# =====================================================================
add_section_divider("03 - Architecture Technique",
                    "Microservices avec API Gateway")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Architecture Technique - Vue d'ensemble", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(5))

# Schema simplifie avec des boites
# Client
client_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.4), Inches(2.8), Inches(0.7))
client_box.fill.solid()
client_box.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
client_box.line.color.rgb = ACCENT_ORANGE
client_box.line.width = Pt(2)
tf = client_box.text_frame
tf.paragraphs[0].text = "Client (Mobile Flutter)"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = ACCENT_ORANGE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.name = "Segoe UI"

# Fleche
add_text_box(slide, Inches(6.1), Inches(2.1), Inches(1.2), Inches(0.4),
             "\u25BC", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Gateway
gw_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(2.5), Inches(5.2), Inches(0.8))
gw_box.fill.solid()
gw_box.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x5B)
gw_box.line.color.rgb = ACCENT_BLUE
gw_box.line.width = Pt(2)
tf = gw_box.text_frame
tf.paragraphs[0].text = "API Gateway (Port 3000) - Routage, Auth, Rate Limiting"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.name = "Segoe UI"

# Fleche
add_text_box(slide, Inches(6.1), Inches(3.3), Inches(1.2), Inches(0.4),
             "\u25BC", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Services - rangee 1
services_row1 = [
    ("Auth\n3001", ACCENT_BLUE),
    ("Booking\n3003", GREEN),
    ("Payment\n3004", ACCENT_ORANGE),
    ("Notification\n3005", YELLOW),
]

for i, (name, color) in enumerate(services_row1):
    x = Inches(0.8) + i * Inches(3.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.8), Inches(2.6), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = name
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Segoe UI"

# Services - rangee 2
services_row2 = [
    ("Restaurant\n3008", RGBColor(0xAF, 0x7A, 0xC5)),
    ("Accommodation\n3009", RGBColor(0x1A, 0xBC, 0x9C)),
    ("Transport\n3011", RGBColor(0xE9, 0x6D, 0x71)),
    ("Service Provider\n3010", RGBColor(0x95, 0xA5, 0xA6)),
]

for i, (name, color) in enumerate(services_row2):
    x = Inches(0.8) + i * Inches(3.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.9), Inches(2.6), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = name
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Segoe UI"

# Resource Core au milieu
rc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(6.0), Inches(4.2), Inches(0.7))
rc_box.fill.solid()
rc_box.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
rc_box.line.color.rgb = LIGHT_GRAY
rc_box.line.width = Pt(1.5)
tf = rc_box.text_frame
tf.paragraphs[0].text = "Resource Core (3002) - Ressources partagees"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.name = "Segoe UI"

# Infra
add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
             "Infrastructure : PostgreSQL (5432) | Redis (6379) | Docker | Kubernetes (prevu)",
             font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 6 - PATTERNS DE COMMUNICATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Patterns de Communication & Securite", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(5))

# Communication patterns
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(3.0),
         "Communication Sync", [
             "HTTP/REST via API Gateway",
             "Swagger/OpenAPI pour doc",
             "Zod pour validation schemas",
         ], accent_color=ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(3.0),
         "Communication Async", [
             "Redis Pub/Sub pour evenements",
             "Queue pour webhooks paiement",
             "Traitement asynchrone",
         ], accent_color=GREEN)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(3.0),
         "Communication Real-time", [
             "WebSocket / Socket.io",
             "Notifications push temps reel",
             "Mise a jour disponibilites live",
         ], accent_color=ACCENT_ORANGE)

# Securite
add_card(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.2),
         "Securite", [
             "JWT avec refresh tokens",
             "Rate limiting (100 req/min par defaut)",
             "CORS configure | Helmet.js pour headers",
             "Authentification centralisee via Auth Service",
         ], accent_color=RED)

add_card(slide, Inches(6.9), Inches(4.8), Inches(5.8), Inches(2.2),
         "Scalabilite", [
             "Chaque service containerisable (Docker)",
             "Orchestration Kubernetes prevue",
             "Cache Redis pour performance",
             "Load balancing au niveau Gateway",
         ], accent_color=RGBColor(0xAF, 0x7A, 0xC5))


# =====================================================================
# SLIDE 7 - LES MICROSERVICES METIER
# =====================================================================
add_section_divider("04 - Les Microservices",
                    "9 services specialises + 1 Gateway")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Services Metier - Restauration", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3.5))

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5),
         "Restaurant Service (Port 3008)", [
             "Creation profil restaurant (nom, type cuisine, capacite)",
             "Gestion des tables par zone (Terrasse, VIP, Interieur...)",
             "Configuration periodes de service (dejeuner, diner...)",
             "Recherche de table disponible pour date/heure/nb personnes",
             "Verification disponibilite (API interne anti-conflit)",
             "Privatisation complete pour evenements",
             "Suggestions creneaux alternatifs",
             "Plan de salle interactif en temps reel",
         ], accent_color=RGBColor(0xAF, 0x7A, 0xC5))

add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.5),
         "Accommodation Service (Port 3009)", [
             "Gestion proprietes et chambres par categorie",
             "Calendrier de disponibilite",
             "Reservation et gestion des sejours",
             "Gestion des avis clients",
             "Tarification dynamique par saison",
             "Check-in / Check-out",
             "Photos et descriptions detaillees",
             "Filtres : prix, equipements, localisation",
         ], accent_color=RGBColor(0x1A, 0xBC, 0x9C))


# =====================================================================
# SLIDE 8 - Transport & Prestations
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Services Metier - Transport & Prestations", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(5))

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5),
         "Transport Service (Port 3011) - LE PLUS AVANCE", [
             "Gestion vehicules et conducteurs",
             "Definition lignes et trajets",
             "Gestion des sieges par vehicule",
             "Reservation de sieges en temps reel",
             "Socket.io pour mises a jour live",
             "Cache Redis pour disponibilites",
             "Prisma ORM pour PostgreSQL",
             "SERVICE ENTIEREMENT IMPLEMENTE",
         ], accent_color=RGBColor(0xE9, 0x6D, 0x71))

add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.5),
         "Service Provider (Port 3010)", [
             "Inscription et gestion des prestataires",
             "Profils marchands detailles",
             "Catalogues de services personnalises",
             "Gestion des horaires et disponibilites",
             "Dashboard marchand (revenus, reservations)",
             "Systeme d'evaluation par les clients",
             "Verification et validation des comptes",
             "Support multi-activites",
         ], accent_color=RGBColor(0x95, 0xA5, 0xA6))


# =====================================================================
# SLIDE 9 - Services Core (Payment, Booking, Auth, Notification)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Services Core - Paiement & Orchestration", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(5))

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.7),
         "Payment Service (3004)", [
             "Stripe (carte bancaire, 3D Secure)",
             "Mobile Money (Orange, MTN, Wave...)",
             "Paiement sur place (especes)",
             "Remboursements & factures PDF",
             "Commissions plateforme",
         ], accent_color=ACCENT_ORANGE)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(2.7),
         "Booking Service (3003)", [
             "Orchestration des reservations",
             "Coordination entre services",
             "Verification disponibilite",
             "Gestion statuts reservation",
             "Historique et annulations",
         ], accent_color=GREEN)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(2.7),
         "Auth Service (3001)", [
             "Inscription / Connexion",
             "JWT + Refresh tokens",
             "OAuth2 (Google, Facebook)",
             "Gestion roles (client/marchand)",
             "OTP verification",
         ], accent_color=ACCENT_BLUE)

add_card(slide, Inches(2.5), Inches(4.5), Inches(3.8), Inches(2.7),
         "Notification Service (3005)", [
             "Emails transactionnels",
             "SMS (confirmation, rappels)",
             "Push notifications mobiles",
             "Templates personnalisables",
             "Envoi evenementiel (Pub/Sub)",
         ], accent_color=YELLOW)

add_card(slide, Inches(6.9), Inches(4.5), Inches(3.8), Inches(2.7),
         "Resource Core (3002)", [
             "Gestion centralisee des ressources",
             "Tables, chambres, sieges, creneaux",
             "API partagee inter-services",
             "Controle de coherence",
             "Allocation multi-ressources",
         ], accent_color=LIGHT_GRAY)


# =====================================================================
# SLIDE 10 - SPECIFICITES CONTEXTE AFRICAIN
# =====================================================================
add_section_divider("05 - Specificites Contexte Africain",
                    "African First - Adapte aux realites locales")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Adapte au Contexte Africain", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(4))

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5),
         "Mobile Money - Priorite absolue", [
             "90% des paiements en Afrique = Mobile Money",
             "Orange Money, MTN Money, Moov Money, Wave, Airtel",
             "Gestion robuste des timeouts et retry (APIs instables)",
             "Confirmation asynchrone (3-5 minutes)",
         ], accent_color=ACCENT_ORANGE)

add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.5),
         "Paiement sur place", [
             "Pratique tres courante (mefiance paiement en ligne)",
             "Option essentielle pour l'adoption de la plateforme",
             "Confirmation manuelle par le commercant",
             "Pas de frais de transaction",
         ], accent_color=GREEN)

add_card(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.5),
         "Devises & Evenements", [
             "Support FCFA (XOF, XAF), USD, EUR",
             "Taux de change dynamiques",
             "Paiement fractionne (acompte + solde)",
             "Cagnottes collectives (mariages, baptemes)",
         ], accent_color=RGBColor(0xAF, 0x7A, 0xC5))

add_card(slide, Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.5),
         "Adaptations UX", [
             "Interface simplifiee pour adoption rapide",
             "Support SMS (pas seulement internet)",
             "Notifications USSD pour Mobile Money",
             "Zones fumeur/non-fumeur (contexte restaurant)",
         ], accent_color=ACCENT_BLUE)


# =====================================================================
# SLIDE 11 - STACK TECHNIQUE
# =====================================================================
add_section_divider("06 - Stack Technique & Outils",
                    "Technologies choisies et justifications")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Stack Technique", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5))

# Backend
add_card(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(3.0),
         "Backend", [
             "TypeScript (typage fort)",
             "Fastify (performance > Express)",
             "Prisma ORM (PostgreSQL)",
             "Redis (cache, sessions, Pub/Sub)",
             "Socket.io (temps reel)",
             "Zod (validation schemas)",
         ], accent_color=ACCENT_BLUE)

# Frontend
add_card(slide, Inches(4.7), Inches(1.4), Inches(3.8), Inches(3.0),
         "Frontend Mobile", [
             "Flutter (Dart) - cross-platform",
             "Dio (requetes HTTP)",
             "Riverpod (gestion d'etat)",
             "26 maquettes SVG realisees",
             "Auth, Client, Marchand flows",
             "Responsive et performant",
         ], accent_color=GREEN)

# Infra
add_card(slide, Inches(8.9), Inches(1.4), Inches(3.8), Inches(3.0),
         "Infrastructure", [
             "Docker & Docker Compose",
             "Kubernetes (prevu)",
             "PostgreSQL (1 DB par service)",
             "Redis (1 cache par service)",
             "Swagger/OpenAPI (docs API)",
             "CI/CD (prevu)",
         ], accent_color=ACCENT_ORANGE)

# Outils dev
add_card(slide, Inches(0.5), Inches(4.7), Inches(5.8), Inches(2.5),
         "Outils de Developpement", [
             "Git avec commits conventionnels (<type>(<scope>): <desc>)",
             "ESLint pour qualite du code",
             "Jest pour les tests unitaires",
             "PlantUML pour diagrammes (cas d'utilisation, classes, sequence)",
             "Pino pour le logging structure",
         ], accent_color=RGBColor(0xAF, 0x7A, 0xC5))

# Securite
add_card(slide, Inches(6.9), Inches(4.7), Inches(5.8), Inches(2.5),
         "Securite & Auth", [
             "JWT + Refresh Tokens",
             "OAuth2 (Google, Facebook)",
             "Helmet.js (headers securises)",
             "Rate Limiting configurable",
             "PCI-DSS compliance (via Stripe)",
         ], accent_color=RED)


# =====================================================================
# SLIDE 12 - MAQUETTES & UX
# =====================================================================
add_section_divider("07 - Maquettes & UX",
                    "26 ecrans concus pour 3 parcours utilisateur")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Maquettes & Parcours Utilisateur", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(4.5))

# Auth flow
add_card(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.5),
         "Parcours Authentification (8 ecrans)", [
             "01 - Splash Screen",
             "02 - Onboarding",
             "03 - Login",
             "04 - Register",
             "05 - OTP Verification",
             "06 - Forgot Password",
             "07 - Reset Password",
             "08 - Inscription Marchand (4 etapes)",
         ], accent_color=ACCENT_BLUE)

# Client flow
add_card(slide, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5),
         "Parcours Client (9+ ecrans)", [
             "01 - Page d'accueil",
             "02 - Recherche etablissement",
             "03 - Detail etablissement",
             "04 - Formulaire reservation",
             "05 - Confirmation reservation",
             "06 - Paiement",
             "07 - Mes reservations",
             "08 - Notifications",
             "09 - Profil utilisateur",
         ], accent_color=GREEN)

# Marchand flow
add_card(slide, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.5),
         "Parcours Marchand", [
             "Dashboard principal",
             "Gestion etablissement",
             "Gestion des reservations",
             "Suivi des paiements",
             "Statistiques d'occupation",
             "Configuration services",
             "Notifications marchands",
             "",
             "Format : SVG vectoriel",
         ], accent_color=ACCENT_ORANGE)


# =====================================================================
# SLIDE 13 - ORGANISATION & METHODOLOGIE
# =====================================================================
add_section_divider("08 - Organisation & Methodologie",
                    "Workflow d'equipe et bonnes pratiques")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Organisation & Methodologie", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(4))

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5),
         "Workflow Git", [
             "Commits conventionnels : <type>(<scope>): <desc>",
             "Branches : feature/, fix/, hotfix/, docs/",
             "Pull Requests obligatoires + code review",
             "Protection branche master (CI doit passer)",
         ], accent_color=ACCENT_BLUE)

add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.5),
         "Documentation", [
             "User Stories par persona (Client, Marchand, System, Admin)",
             "Diagrammes PlantUML (cas d'utilisation, classes, sequence)",
             "Architecture technique documentee",
             "Guide de contribution (CONTRIBUTING.md)",
         ], accent_color=GREEN)

add_card(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.5),
         "Gestion de projet Agile", [
             "User Stories priorisees (Haute / Moyenne / Basse)",
             "Estimation en story points (Fibonacci)",
             "Roadmap par sprint (2 semaines/sprint)",
             "Equipe de 4 developpeurs + Tech Lead",
         ], accent_color=ACCENT_ORANGE)

add_card(slide, Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.5),
         "Templates & Standards", [
             "Template de PR (checklist, description, tests)",
             "Templates d'issues (bug report, feature request)",
             "Git hooks recommandes (pre-commit)",
             "ESLint + standards de code documentes",
         ], accent_color=RGBColor(0xAF, 0x7A, 0xC5))


# =====================================================================
# SLIDE 14 - DIAGRAMMES UML
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Documentation UML - Diagrammes", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(4))

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Chaque microservice dispose de ses propres diagrammes PlantUML :",
             font_size=18, color=LIGHT_GRAY)

services_diag = [
    "Auth Service", "Booking Service", "Restaurant Service",
    "Accommodation Service", "Transport Service", "Payment Service",
    "Notification Service", "Resource Core", "Service Provider"
]

for i, svc in enumerate(services_diag):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(2.2) + row * Inches(1.6)
    add_card(slide, x, y, Inches(3.8), Inches(1.4), svc, [
        "Cas d'utilisation",
        "Diagramme de classes",
        "Diagramme de sequence",
    ], accent_color=ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "+ Diagrammes globaux : deploiement, interoperabilite, architecture Excalidraw",
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 15 - ETAT D'AVANCEMENT
# =====================================================================
add_section_divider("09 - Etat d'avancement",
                    "Ou en sommes-nous ?")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Etat d'avancement", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3))

# Fait
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.5),
         "TERMINE", [
             "Architecture globale definie",
             "API Gateway configure",
             "User Stories documentees (tous services)",
             "Maquettes UI/UX (26 ecrans SVG)",
             "Transport Service IMPLEMENTE",
             "Diagrammes UML (tous services)",
             "Guide de contribution",
             "Infrastructure Docker",
         ], accent_color=GREEN)

# En cours
add_card(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(5.5),
         "EN COURS", [
             "Auth Service (implementation)",
             "Booking Service (implementation)",
             "Payment Service (implementation)",
             "Restaurant Service (implementation)",
             "Accommodation Service",
             "Notification Service",
             "Application mobile Flutter",
             "",
         ], accent_color=YELLOW)

# A venir
add_card(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(5.5),
         "A VENIR", [
             "Deploiement Kubernetes",
             "Pipeline CI/CD",
             "Tests unitaires complets",
             "Tests d'integration",
             "Documentation API Swagger",
             "Beta testing",
             "Mise en production",
             "",
         ], accent_color=RED)


# =====================================================================
# SLIDE 16 - CHIFFRES CLES
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Chiffres Cles du Projet", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3))

metrics = [
    ("9+", "Microservices", ACCENT_BLUE),
    ("26", "Maquettes UI", GREEN),
    ("17", "User Stories\nPayment", ACCENT_ORANGE),
    ("15", "User Stories\nRestaurant", RGBColor(0xAF, 0x7A, 0xC5)),
    ("5", "Moyens de\npaiement", YELLOW),
    ("4", "Devs dans\nl'equipe", RGBColor(0x1A, 0xBC, 0x9C)),
]

for i, (number, label, color) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(1.0) + col * Inches(4.0)
    y = Inches(1.8) + row * Inches(2.8)

    # Number
    add_text_box(slide, x, y, Inches(3.2), Inches(1.2),
                 number, font_size=60, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, x, y + Inches(1.2), Inches(3.2), Inches(0.8),
                 label, font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 17 - DEMO & PERSPECTIVES
# =====================================================================
add_section_divider("10 - Demo & Perspectives",
                    "Prochaines etapes et vision long terme")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Perspectives & Prochaines Etapes", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(4.5))

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.3),
         "Court terme (1-2 mois)", [
             "Finaliser l'implementation de tous les services backend",
             "Developper l'application mobile Flutter",
             "Mettre en place les tests unitaires et d'integration",
             "Deployer un environnement de staging",
         ], accent_color=GREEN)

add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.3),
         "Moyen terme (3-6 mois)", [
             "Integration reelle Mobile Money (APIs operateurs)",
             "Beta testing avec des commercants locaux",
             "Pipeline CI/CD complet",
             "Deploiement Kubernetes en production",
         ], accent_color=ACCENT_ORANGE)

add_card(slide, Inches(0.5), Inches(4.1), Inches(5.8), Inches(2.3),
         "Long terme (6-12 mois)", [
             "Intelligence artificielle : suggestions personnalisees",
             "Analytics avances pour les commercants",
             "Expansion multi-pays (devises, operateurs)",
             "Marketplace de services complementaires",
         ], accent_color=RGBColor(0xAF, 0x7A, 0xC5))

add_card(slide, Inches(6.9), Inches(4.1), Inches(5.8), Inches(2.3),
         "Vision", [
             "Devenir LA plateforme de reference de reservation",
             "en Afrique, adaptee aux realites du terrain",
             "Multi-services, multi-devises, multi-paiements",
             "Accessibilite maximale (SMS, USSD, Mobile Money)",
         ], accent_color=ACCENT_BLUE)


# =====================================================================
# SLIDE 18 - CONCLUSION / MERCI
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide, SECTION_BG)

add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(0.08), W, Inches(0.04), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.0),
             "Merci de votre attention", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(3.0), Inches(4.3), ACCENT_ORANGE)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.0),
             "BookingSysteme - Plateforme de reservation multi-services\nAdaptee au contexte africain",
             font_size=22, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.6),
             "Equipe PING 8 | Fevrier 2026",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.8),
             "Questions ?",
             font_size=36, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SAUVEGARDE
# =====================================================================
output_path = os.path.join(os.path.dirname(__file__), "BookingSysteme_Presentation.pptx")
prs.save(output_path)
print(f"Presentation generee avec succes : {output_path}")
print(f"Nombre de slides : {len(prs.slides)}")
